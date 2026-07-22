# -*- coding: utf-8 -*-
"""Generates the UTMS Enterprise Automation Roadmap presentation (.pptx)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------- palette --
def rgb(h):
    return RGBColor((h >> 16) & 255, (h >> 8) & 255, h & 255)

NAVY = 0x0B2A4A
NAVY_DARK = 0x081E36
ACCENT = 0x1F6FEB
ACCENT_SOFT = 0x6FA0F5
ACCENT_PALE = 0xBFD6FA
LIGHT_BLUE = 0xEDF3FC
LIGHT_BLUE2 = 0xE3ECFB
CARD_BORDER = 0xCFDFF5
WHITE = 0xFFFFFF
TEXT_DARK = 0x18283F
TEXT_GRAY = 0x60708A
TEXT_LIGHT_GRAY = 0x8B98AC
RED = 0xB6362C
RED_BG = 0xFBEAE8
AMBER = 0xAD6F0B
AMBER_BG = 0xFCF1DD
GREEN = 0x1E7A3B
GREEN_BG = 0xE6F4EA
GRAY_LINE = 0xE2E7F0

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
CL = Inches(0.55)   # content left
CR = Inches(12.78)  # content right
CW = CR - CL        # content width

FONT = "Calibri"


def _e(v):
    """Coerce any numeric (int, Length, or stray float from true-division) to a
    proper integer-EMU Length so python-pptx never writes a float coordinate
    (e.g. '0.0') into the XML -- PowerPoint's strict parser rejects those even
    though python-pptx itself reads them back without complaint."""
    return Emu(int(round(v)))


# ------------------------------------------------------------- primitives --
def add_shadow(shape, blur=38000, dist=18000, direction=5400000, color=0x9AA7B8, alpha=32000):
    spPr = shape._element.spPr
    existing = spPr.find(qn('a:effectLst'))
    if existing is not None:
        spPr.remove(existing)
    eff = etree.SubElement(spPr, qn('a:effectLst'))
    shdw = etree.SubElement(eff, qn('a:outerShdw'))
    shdw.set('blurRad', str(blur))
    shdw.set('dist', str(dist))
    shdw.set('dir', str(direction))
    shdw.set('rotWithShape', '0')
    clr = etree.SubElement(shdw, qn('a:srgbClr'))
    clr.set('val', '%06X' % color)
    a = etree.SubElement(clr, qn('a:alpha'))
    a.set('val', str(alpha))


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_w=1.0,
             shadow=False, shape_type=MSO_SHAPE.RECTANGLE, radius=None):
    shp = slide.shapes.add_shape(shape_type, _e(x), _e(y), _e(w), _e(h))
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
    else:
        shp.fill.background()
    if line_color is not None:
        shp.line.color.rgb = rgb(line_color)
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if radius is not None and shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        shp.adjustments[0] = radius
    if shadow:
        add_shadow(shp)
    return shp


def _set_run(p, text, size, bold, color, font=FONT, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = rgb(color)
    r.font.name = font
    return r


def add_text(slide, x, y, w, h, text, size=10, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, font=FONT, italic=False, anchor=MSO_ANCHOR.TOP,
             wrap=True, line_spacing=None):
    box = slide.shapes.add_textbox(_e(x), _e(y), _e(w), _e(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    _set_run(p, lines[0], size, bold, color, font, italic)
    for extra in lines[1:]:
        p2 = tf.add_paragraph()
        p2.alignment = align
        if line_spacing:
            p2.line_spacing = line_spacing
        _set_run(p2, extra, size, bold, color, font, italic)
    return box


def fill_shape_text(shape, text, size=10, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, font=FONT, wrap=True, anchor=MSO_ANCHOR.MIDDLE,
                     margins=0.02):
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(margins)
    tf.margin_right = Inches(margins)
    tf.margin_top = Inches(margins)
    tf.margin_bottom = Inches(margins)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    _set_run(p, text, size, bold, color, font)
    return shape


def add_icon_badge(slide, x, y, size, label, fill=ACCENT, text_color=WHITE, font_size=8.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _e(x), _e(y), _e(size), _e(size))
    shp.adjustments[0] = 0.28
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.fill.background()
    shp.shadow.inherit = False
    fill_shape_text(shp, label, size=font_size, bold=True, color=text_color)
    return shp


def add_connector(slide, x1, y1, x2, y2, color=ACCENT, width=1.25, dashed=False,
                   arrow=True, arrow_type='triangle'):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _e(x1), _e(y1), _e(x2), _e(y2))
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    if dashed:
        d = etree.SubElement(ln, qn('a:prstDash'))
        d.set('val', 'dash')
    if arrow:
        te = etree.SubElement(ln, qn('a:tailEnd'))
        te.set('type', arrow_type)
        te.set('w', 'med')
        te.set('len', 'med')
    return conn


def decision_diamond(slide, x, y, w, h, text, fill=LIGHT_BLUE, border=ACCENT,
                      text_color=NAVY, size=9.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, _e(x), _e(y), _e(w), _e(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(border)
    shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    add_shadow(shp)
    fill_shape_text(shp, text, size=size, bold=True, color=text_color, margins=0.12)
    return shp


def process_box(slide, x, y, w, h, title, desc, icon_label, fill=WHITE,
                 border=ACCENT, title_color=NAVY, desc_color=TEXT_GRAY,
                 icon_fill=ACCENT, title_size=10.5, desc_size=8, shadow=True,
                 badge_size=None):
    if badge_size is None:
        badge_size = Inches(0.34)
    shp = add_rect(slide, x, y, w, h, fill=fill, line_color=border, line_w=1.25,
                   shadow=shadow, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09)
    pad = Inches(0.09)
    bs = min(badge_size, h - 2 * pad) if h - 2 * pad > 0 else badge_size
    badge_y = y + (h - bs) / 2
    add_icon_badge(slide, x + pad, badge_y, bs, icon_label, fill=icon_fill, font_size=8)
    text_x = x + pad + bs + Inches(0.09)
    text_w = w - pad - bs - Inches(0.09) - pad
    box = slide.shapes.add_textbox(_e(text_x), _e(y + pad * 0.35), _e(text_w), _e(h - pad * 0.7))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p0 = tf.paragraphs[0]
    _set_run(p0, title, title_size, True, title_color)
    if desc:
        p1 = tf.add_paragraph()
        p1.space_before = Pt(1)
        p1.line_spacing = 1.02
        _set_run(p1, desc, desc_size, False, desc_color)
    return shp


def compact_box(slide, x, y, w, h, title, caption, icon_label, fill=WHITE,
                 border=ACCENT, icon_fill=ACCENT, title_size=9, caption_size=7.2):
    shp = add_rect(slide, x, y, w, h, fill=fill, line_color=border, line_w=1.1,
                   shadow=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    bs = Inches(0.24)
    top_pad = Inches(0.07)
    bx = x + (w - bs) / 2
    add_icon_badge(slide, bx, y + top_pad, bs, icon_label, fill=icon_fill, font_size=7)
    text_y = y + top_pad + bs + Inches(0.03)
    text_h = h - (text_y - y) - Inches(0.05)
    box = slide.shapes.add_textbox(_e(x + Inches(0.04)), _e(text_y), _e(w - Inches(0.08)), _e(text_h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    _set_run(p0, title, title_size, True, NAVY)
    if caption:
        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        p1.line_spacing = 1.0
        _set_run(p1, caption, caption_size, False, TEXT_GRAY)
    return shp


def callout_card(slide, x, y, w, h, title, desc, bar_color, bg_tint, title_size=11, desc_size=8.5):
    shp = add_rect(slide, x, y, w, h, fill=bg_tint, line_color=None,
                   shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08, shadow=True)
    add_rect(slide, x, y, Inches(0.09), h, fill=bar_color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    pad = Inches(0.16)
    add_text(slide, x + pad + Inches(0.06), y + Inches(0.1), w - pad - Inches(0.2), Inches(0.28),
              title, size=title_size, bold=True, color=bar_color)
    add_text(slide, x + pad + Inches(0.06), y + Inches(0.42), w - pad - Inches(0.2), h - Inches(0.5),
              desc, size=desc_size, bold=False, color=TEXT_DARK, line_spacing=1.05)
    return shp


def pain_row(slide, x, y, w, title, desc, code, title_size=11.5, desc_size=9.5):
    bs = Inches(0.42)
    add_icon_badge(slide, x, y, bs, code, fill=ACCENT, font_size=9, text_color=WHITE)
    add_text(slide, x + bs + Inches(0.16), y - Inches(0.02), w - bs - Inches(0.16), Inches(0.28),
              title, size=title_size, bold=True, color=NAVY)
    add_text(slide, x + bs + Inches(0.16), y + Inches(0.26), w - bs - Inches(0.16), Inches(0.55),
              desc, size=desc_size, bold=False, color=TEXT_GRAY, line_spacing=1.05)


def challenge_card(slide, x, y, w, h, title, desc, code, title_size=11, desc_size=8.5):
    shp = add_rect(slide, x, y, w, h, fill=WHITE, line_color=CARD_BORDER, line_w=1,
                   shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07, shadow=True)
    add_rect(slide, x, y, w, Inches(0.07), fill=ACCENT, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    bs = Inches(0.36)
    add_icon_badge(slide, x + Inches(0.16), y + Inches(0.18), bs, code, fill=NAVY, font_size=8)
    add_text(slide, x + Inches(0.16), y + Inches(0.18) + bs + Inches(0.08), w - Inches(0.32), Inches(0.5),
              title, size=title_size, bold=True, color=NAVY, line_spacing=1.0)
    add_text(slide, x + Inches(0.16), y + h - (h - (Inches(0.18) + bs + Inches(0.08) + Inches(0.45))), w - Inches(0.32),
              h - (Inches(0.18) + bs + Inches(0.08) + Inches(0.45)) - Inches(0.12),
              desc, size=desc_size, bold=False, color=TEXT_GRAY, line_spacing=1.06)
    return shp


def add_header(slide, kicker, title, title_size=25, badge_text=None):
    add_rect(slide, 0, 0, Inches(0.14), SLIDE_H, fill=NAVY)
    add_text(slide, CL, Inches(0.3), Inches(10.5), Inches(0.26), kicker.upper(),
              size=11.5, bold=True, color=ACCENT)
    add_text(slide, CL, Inches(0.56), Inches(11.9), Inches(0.55), title,
              size=title_size, bold=True, color=NAVY)
    add_rect(slide, CL, Inches(1.145), Inches(1.15), Pt(3), fill=ACCENT)
    add_rect(slide, CL, Inches(1.24), CW, Pt(0.75), fill=GRAY_LINE)
    if badge_text:
        bw = Inches(0.14) + Inches(0.082) * len(badge_text)
        bx = CR - bw
        add_rect(slide, bx, Inches(0.32), bw, Inches(0.3), fill=NAVY,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        add_text(slide, bx, Inches(0.32), bw, Inches(0.3), badge_text, size=8.5, bold=True,
                  color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, page_no, total=7):
    add_rect(slide, CL, Inches(7.08), CW, Pt(0.75), fill=GRAY_LINE)
    add_text(slide, CL, Inches(7.15), Inches(7), Inches(0.28),
              "UTMS Enterprise Automation Program  ·  Confidential", size=8.5, color=TEXT_LIGHT_GRAY)
    add_text(slide, Inches(11.6), Inches(7.15), Inches(1.18), Inches(0.28),
              f"{page_no} / {total}", size=8.5, color=TEXT_LIGHT_GRAY, align=PP_ALIGN.RIGHT)


def legend_strip(slide, x, y):
    items = [
        (ACCENT, MSO_SHAPE.ROUNDED_RECTANGLE, "Automated processing step"),
        (ACCENT, MSO_SHAPE.DIAMOND, "Decision point"),
        (None, None, "Cross-cutting service (dashed link)"),
    ]
    cx = x
    for color, shape, label in items:
        if shape == MSO_SHAPE.DIAMOND:
            sw = Inches(0.16)
            sh = Inches(0.16)
            s = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, _e(cx), _e(y), _e(sw), _e(sh))
            s.fill.solid(); s.fill.fore_color.rgb = rgb(LIGHT_BLUE)
            s.line.color.rgb = rgb(ACCENT); s.line.width = Pt(1)
            s.shadow.inherit = False
            cx2 = cx + sw + Inches(0.08)
        elif shape == MSO_SHAPE.ROUNDED_RECTANGLE:
            sw = Inches(0.22)
            sh = Inches(0.13)
            add_rect(slide, cx, y + Inches(0.015), sw, sh, fill=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
            cx2 = cx + sw + Inches(0.08)
        else:
            sw = Inches(0.22)
            add_connector(slide, cx, y + Inches(0.08), cx + sw, y + Inches(0.08), color=ACCENT, width=1.25, dashed=True, arrow=False)
            cx2 = cx + sw + Inches(0.08)
        add_text(slide, cx2, y - Inches(0.02), Inches(2.6), Inches(0.2), label, size=8, color=TEXT_GRAY)
        cx = cx2 + Inches(2.65)


NOTES = {}


def set_notes(slide, idx, text):
    ns = slide.notes_slide
    ns.notes_text_frame.text = text


# ================================================================ BUILD ===
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]
TOTAL_SLIDES = 7

# ---------------------------------------------------------------- Slide 0 --
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
add_rect(s, 0, 0, Inches(5.3), SLIDE_H, fill=NAVY)
# decorative circles on navy panel
c1 = add_rect(s, Inches(-0.8), Inches(5.6), Inches(2.6), Inches(2.6), fill=NAVY_DARK, shape_type=MSO_SHAPE.OVAL)
c2 = add_rect(s, Inches(3.6), Inches(-1.0), Inches(2.2), Inches(2.2), fill=0x123A63, shape_type=MSO_SHAPE.OVAL)
add_rect(s, Inches(0.6), Inches(1.72), Inches(0.9), Pt(3.5), fill=ACCENT_SOFT)
add_text(s, Inches(0.6), Inches(1.9), Inches(4.3), Inches(0.3), "ENTERPRISE AUTOMATION BLUEPRINT",
         size=12, bold=True, color=ACCENT_SOFT)
add_text(s, Inches(0.6), Inches(2.28), Inches(4.35), Inches(1.55), "Enterprise Automation\nRoadmap",
         size=33, bold=True, color=WHITE, line_spacing=1.02)
add_text(s, Inches(0.6), Inches(3.92), Inches(4.35), Inches(0.85),
         "Unified Ticket Management System — Healthcare Revenue Cycle Management",
         size=13.5, bold=False, color=ACCENT_PALE, line_spacing=1.15)
add_rect(s, Inches(0.6), Inches(4.92), Inches(0.9), Pt(2.5), fill=ACCENT)
add_text(s, Inches(0.6), Inches(5.1), Inches(4.3), Inches(0.24), "PREPARED FOR", size=9.5, bold=True, color=ACCENT_SOFT)
add_text(s, Inches(0.6), Inches(5.36), Inches(4.35), Inches(0.55),
         "Chief Technology Officer · Enterprise Architecture\nProduct & Client Leadership", size=11, color=WHITE, line_spacing=1.15)
add_text(s, Inches(0.6), Inches(6.9), Inches(4.4), Inches(0.35),
         "Confidential — For internal and client discussion", size=9, italic=True, color=0x7C97BE)

# right side preview cards
add_text(s, Inches(6.0), Inches(1.55), Inches(6.6), Inches(0.3), "TWO PRIORITY AUTOMATION USE CASES",
         size=11, bold=True, color=ACCENT)
card1 = add_rect(s, Inches(6.0), Inches(1.95), Inches(6.65), Inches(2.0), fill=WHITE, line_color=CARD_BORDER,
                 line_w=1, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, shadow=True)
add_icon_badge(s, Inches(6.24), Inches(2.2), Inches(0.5), "01", fill=NAVY, font_size=13)
add_text(s, Inches(6.94), Inches(2.2), Inches(5.5), Inches(0.35), "Joiner–Mover–Leaver (JML) Lifecycle Automation",
         size=15, bold=True, color=NAVY)
add_text(s, Inches(6.94), Inches(2.62), Inches(5.5), Inches(0.75),
         "End-to-end identity, access, and ticket-ownership automation across the full employee lifecycle — hire to exit.",
         size=10.5, color=TEXT_GRAY, line_spacing=1.15)
add_text(s, Inches(6.24), Inches(3.55), Inches(6.2), Inches(0.3),
         "RBAC Automation · Identity Governance · Organization Sync", size=9.5, bold=True, color=ACCENT)

card2 = add_rect(s, Inches(6.0), Inches(4.15), Inches(6.65), Inches(2.0), fill=WHITE, line_color=CARD_BORDER,
                 line_w=1, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, shadow=True)
add_icon_badge(s, Inches(6.24), Inches(4.4), Inches(0.5), "02", fill=NAVY, font_size=13)
add_text(s, Inches(6.94), Inches(4.4), Inches(5.5), Inches(0.35), "Intelligent Ticket Lifecycle Automation",
         size=15, bold=True, color=NAVY)
add_text(s, Inches(6.94), Inches(4.82), Inches(5.5), Inches(0.75),
         "Rules-based intake, triage, and status tracking — deterministic, auditable, and cost-effective, with the Account Manager always holding final control.",
         size=10.5, color=TEXT_GRAY, line_spacing=1.15)
add_text(s, Inches(6.24), Inches(5.75), Inches(6.2), Inches(0.3),
         "Rules-Based Intake · Human-in-the-Loop · Cost-Effective", size=9.5, bold=True, color=ACCENT)

add_text(s, Inches(6.0), Inches(6.95), Inches(6.6), Inches(0.3),
         "Automation patterns adapted from ServiceNow · Okta · Salesforce · Workday · Freshservice",
         size=9, italic=True, color=TEXT_LIGHT_GRAY)

set_notes(s, 0, "Title slide. Frame the deck: two enterprise-grade automation use cases for UTMS, benchmarked against "
                 "mature patterns from ServiceNow, Okta, Salesforce, Workday, and Freshservice, then adapted specifically "
                 "to UTMS's current architecture. Use Case 1 addresses the employee lifecycle (Joiner-Mover-Leaver); "
                 "Use Case 2 addresses AI-assisted ticket intake with the Account Manager retaining final authority. "
                 "Emphasize this is a roadmap for CTO/architecture/product leadership, not a finished build.")

# ---------------------------------------------------------------- Slide 1 --
# UC1 - Pain Points
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
add_header(s, "Use Case 1 · Pain Points", "Joiner–Mover–Leaver (JML) Lifecycle Automation")

left_x = CL
left_w = Inches(6.9)
add_text(s, left_x, Inches(1.42), left_w, Inches(0.3), "Current Manual Workflow", size=13.5, bold=True, color=NAVY)

pain_items = [
    ("JN", "Joiner Onboarding",
     "Super Admin manually creates the account, role, category, and hierarchy links field-by-field with no validation, whenever a new joiner is identified."),
    ("MV", "Mover / Transfer",
     "Promotions and department changes require manually revoking old permissions and re-granting new ones across screens."),
    ("LV", "Leaver Offboarding",
     "Deactivation depends on someone remembering to disable the account, revoke access, and reassign open tickets."),
    ("AC", "Approval & Hierarchy",
     "Manager/team-lead consistency is checked at save time, but nothing orchestrates the lifecycle end-to-end."),
    ("AU", "Audit Reconstruction",
     "Lifecycle history is pieced together after the fact from scattered logs rather than captured as one coherent trail."),
]
row_y = 1.84
for code, title, desc in pain_items:
    pain_row(s, left_x, Inches(row_y), left_w, title, desc, code)
    row_y += 0.95

right_x = Inches(7.78)
right_w = Inches(5.0)
add_text(s, right_x, Inches(1.42), right_w, Inches(0.3), "Business Impact", size=13.5, bold=True, color=NAVY)
impacts = [
    ("Security Risk", "Orphaned access after offboarding is the #1 finding in internal access audits.", RED, RED_BG),
    ("Compliance Exposure", "No continuous evidence trail to support HIPAA minimum-necessary-access reviews.", AMBER, AMBER_BG),
    ("Operational Delay", "Multi-day gap between hire date and full productivity while access is assembled by hand.", ACCENT, LIGHT_BLUE),
    ("Cost Impact", "Recurring manual reconciliation consumes significant admin and compliance hours every cycle.", NAVY, LIGHT_BLUE2),
]
cy = 1.84
for title, desc, bar, bg in impacts:
    callout_card(s, right_x, Inches(cy), right_w, Inches(1.06), title, desc, bar, bg)
    cy += 1.19

add_footer(s, 2, TOTAL_SLIDES)
set_notes(s, 1,
    "Walk through today's fully manual JML process: there is no HR system integration -- Super Admin is the sole "
    "source of truth for who's joining, moving, or leaving, and provisions each joiner field-by-field with no validation; "
    "movers require manual revoke-then-grant across multiple screens; leavers depend on someone remembering "
    "to deactivate and reassign. Land on the four business-impact callouts: orphaned access is the top insider-risk "
    "finding in access audits, there's no continuous HIPAA minimum-necessary evidence trail, onboarding-to-productivity "
    "takes multiple days, and manual reconciliation is a recurring cost every audit cycle.")

# ---------------------------------------------------------------- Slide 2 --
# UC1 - Solution Architecture
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
add_header(s, "Use Case 1 · Solution Architecture", "Joiner–Mover–Leaver (JML) Lifecycle Automation", title_size=23)

CX = Inches(3.15)
CWID = Inches(6.7)
RAIL_X = Inches(10.05)
RAIL_W = Inches(2.73)
LRAIL_X = Inches(0.55)
LRAIL_W = Inches(2.4)

# Row 1 - Inputs
bw = Inches((6.7 - 0.2) / 2)
process_box(s, CX, Inches(1.42), bw, Inches(0.55), "Super Admin", "Identifies and initiates a joiner / mover / leaver action via the Admin Console — no external HR feed.", "SA", icon_fill=NAVY, badge_size=Inches(0.3), title_size=9.5, desc_size=7.3)
process_box(s, CX + bw + Inches(0.2), Inches(1.42), bw, Inches(0.55), "Identity Layer", "Authenticates the actor and issues the session / JWT context.", "ID", icon_fill=NAVY, badge_size=Inches(0.3), title_size=9.5, desc_size=7.3)
add_connector(s, CX + CWID / 2, Inches(1.97), CX + CWID / 2, Inches(2.14), color=ACCENT)
add_text(s, CX + CWID / 2 + Inches(0.08), Inches(1.955), Inches(2.0), Inches(0.2), "Lifecycle Action", size=7.5, italic=True, color=TEXT_GRAY)

# Row 2 - Automation trigger
process_box(s, CX, Inches(2.16), CWID, Inches(0.5), "Employee Event Engine",
            "Captures the Super Admin-initiated action and starts the correct lifecycle workflow.", "EVT", badge_size=Inches(0.3), title_size=10, desc_size=7.5)
add_connector(s, CX + CWID / 2, Inches(2.66), CX + CWID / 2, Inches(2.81), color=ACCENT)
add_text(s, CX + CWID / 2 + Inches(0.08), Inches(2.645), Inches(2.4), Inches(0.2), "Validated Change Set", size=7.5, italic=True, color=TEXT_GRAY)

# Row 3 - Processing services (4 compact boxes)
n = 4
gap = Inches(0.14)
cbw = Emu(int((CWID - gap * (n - 1)) / n))
cbx = CX
proc_items = [
    ("RB", "RBAC Engine", "role defaults"),
    ("PR", "Permission Resolver", "effective delta"),
    ("OR", "Organization Svc", "hierarchy update"),
    ("WF", "Workflow Engine", "task sequencing"),
]
for code, title, cap in proc_items:
    compact_box(s, cbx, Inches(2.83), cbw, Inches(0.66), title, cap, code)
    cbx = Emu(int(cbx + cbw + gap))
add_connector(s, CX + CWID / 2, Inches(3.49), CX + CWID / 2, Inches(3.63), color=ACCENT)
add_text(s, CX + CWID / 2 + Inches(0.08), Inches(3.475), Inches(1.6), Inches(0.2), "Approved", size=7.5, italic=True, color=TEXT_GRAY)

# Row 4 - Decision diamond
dw = Inches(2.7)
dh = Inches(0.64)
dx = CX + (CWID - dw) / 2
decision_diamond(s, dx, Inches(3.65), dw, dh, "Consistency & Approval Check", size=9)
add_connector(s, dx + dw / 2, Inches(4.29), dx + dw / 2, Inches(4.42), color=ACCENT)

# left dashed connector to exception rail
add_connector(s, dx, Inches(3.65) + dh / 2, LRAIL_X + LRAIL_W, Inches(3.65) + dh / 2, color=AMBER, dashed=True, width=1.1)
add_text(s, LRAIL_X + LRAIL_W - Inches(1.55), Inches(3.65) + dh / 2 - Inches(0.24), Inches(1.5), Inches(0.2),
         "Conflict / Rejected", size=7, italic=True, color=AMBER, align=PP_ALIGN.RIGHT)

# Row 5 - Business / provisioning services (3 boxes)
n = 3
gap = Inches(0.16)
cbw2 = Emu(int((CWID - gap * (n - 1)) / n))
cbx = CX
biz_items = [
    ("PV", "Account Provisioning", "creates account & dashboards"),
    ("XF", "Ticket & Ownership Transfer", "reassigns tickets & reports"),
    ("SP", "Compliance Snapshot", "packages change for audit"),
]
for code, title, cap in biz_items:
    compact_box(s, cbx, Inches(4.44), cbw2, Inches(0.66), title, cap, code, title_size=8.3, caption_size=6.8)
    cbx = Emu(int(cbx + cbw2 + gap))
add_connector(s, CX + CWID / 2, Inches(5.10), CX + CWID / 2, Inches(5.24), color=ACCENT)

# Row 6 - Outputs
bw = Inches((6.7 - 0.2) / 2)
process_box(s, CX, Inches(5.26), bw, Inches(0.56), "User Portal",
            "Employee profile & dashboard reflect the new state immediately.", "PT", icon_fill=NAVY, badge_size=Inches(0.3), title_size=9.5, desc_size=7.3)
process_box(s, CX + bw + Inches(0.2), Inches(5.26), bw, Inches(0.56), "Management Dashboard",
            "Supervisors see org / permission changes in real time.", "MG", icon_fill=NAVY, badge_size=Inches(0.3), title_size=9.5, desc_size=7.3)

# Right rail - cross cutting services
rail_items = [
    ("NT", "Notification Service", "Alerts employee, manager, and approver at every lifecycle step."),
    ("AUD", "Audit Service", "Writes an immutable event for every automated grant, revoke, and transfer."),
    ("RPT", "Reporting Service", "Feeds compliance and access-recertification reporting."),
]
rail_y = [2.16, 3.32, 4.48]
rail_h = Inches(0.98)
for (code, title, desc), ry in zip(rail_items, rail_y):
    process_box(s, RAIL_X, Inches(ry), RAIL_W, rail_h, title, desc, code, fill=LIGHT_BLUE2, border=ACCENT_SOFT,
                icon_fill=ACCENT, badge_size=Inches(0.3), title_size=9, desc_size=7.3)
for i in range(len(rail_y) - 1):
    add_connector(s, RAIL_X + RAIL_W / 2, Inches(rail_y[i]) + rail_h, RAIL_X + RAIL_W / 2, Inches(rail_y[i + 1]), color=ACCENT_SOFT, width=1)
# dashed cross-links from central flow into the rail
add_connector(s, CX + CWID, Inches(2.41), RAIL_X, Inches(rail_y[0]) + rail_h / 2, color=ACCENT_SOFT, dashed=True, width=1, arrow=False)
add_connector(s, CX + CWID, Inches(3.16), RAIL_X, Inches(rail_y[1]) + rail_h / 2, color=ACCENT_SOFT, dashed=True, width=1, arrow=False)
add_connector(s, CX + CWID, Inches(4.77), RAIL_X, Inches(rail_y[2]) + rail_h / 2, color=ACCENT_SOFT, dashed=True, width=1, arrow=False)

# Left rail - exception handling
callout_card(s, LRAIL_X, Inches(3.98), LRAIL_W, Inches(0.98),
             "Exception Handling",
             "A failed validation pauses the workflow and routes to Admin for manual resolution — nothing fails silently.",
             AMBER, AMBER_BG, title_size=10, desc_size=8)

legend_strip(s, CL, Inches(6.05))
add_footer(s, 3, TOTAL_SLIDES)
set_notes(s, 2,
    "This is the architecture slide. There's no external HR system in this design -- Super Admin, who already tracks "
    "who's joining, moving, or leaving, initiates the action directly via the Admin Console; Identity Layer supplies the "
    "authenticated session context. Both feed the Employee Event Engine, which classifies the action as Joiner, Mover, or "
    "Leaver. Four processing services compute the change: RBAC Engine (role defaults), "
    "Permission Resolver (effective delta), Organization Service (hierarchy update), Workflow Engine (task sequencing). "
    "A Consistency & Approval Check gate either proceeds to provisioning/transfer/compliance-snapshot services or routes "
    "to Exception Handling on the left for manual resolution -- nothing fails silently. Outputs land in the User Portal "
    "and Management Dashboard. Notification, Audit, and Reporting run as cross-cutting services on the right rail, "
    "receiving events from every stage of the flow -- this is the same pattern ServiceNow uses for Flow Designer "
    "orchestration with parallel notification/audit logging.")

# ---------------------------------------------------------------- Slide 3 --
# UC1 - Implementation Challenges
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
add_header(s, "Use Case 1 · Implementation Challenges", "Joiner–Mover–Leaver (JML) Lifecycle Automation")

challenges_uc1 = [
    ("SYN", "Identity Synchronization", "Keeping a Super-Admin-triggered action consistent with the affected user's live session state, without double-processing a resubmitted action."),
    ("TRG", "Manual Trigger Reliance", "With no external HR feed, the automation is only as timely as Super Admin noticing and initiating the action — a missed trigger means the workflow never starts."),
    ("CNF", "Permission Conflicts", "The additive-only override model must be extended carefully to avoid unintended privilege stacking."),
    ("INH", "Role Inheritance", "Hierarchy changes must cascade correctly without breaking existing scoped or ticket-level grants."),
    ("SCL", "Scalability", "The lifecycle sweep and cache-invalidation path must hold up as headcount and org depth grow."),
    ("RBK", "Rollback Strategy", "A failed Mover / Leaver step must be safely reversible — never left half-applied."),
    ("SEC", "Security & Compliance", "Every automated grant or revoke must produce HIPAA-defensible, immutable evidence."),
    ("OPT", "Analytics Enhancement (Optional)", "Anomaly detection on lifecycle patterns is a future, data-justified add-on — the core workflow needs no AI/LLM dependency at all."),
]
cols, rows = 4, 2
gap = Inches(0.2)
card_w = Emu(int((CW - gap * (cols - 1)) / cols))
card_h = Inches(2.35)
y0 = 1.55
row_gap = 0.35
for i, (code, title, desc) in enumerate(challenges_uc1):
    r, c = divmod(i, cols)
    x = Emu(int(CL + c * (card_w + gap)))
    y = Inches(y0 + r * (2.35 + row_gap))
    challenge_card(s, x, y, card_w, card_h, title, desc, code)

add_footer(s, 4, TOTAL_SLIDES)
set_notes(s, 3,
    "Key implementation risks for JML automation: identity sync for a Super-Admin-triggered action, the fact that the "
    "whole workflow now depends on Super Admin noticing and manually triggering each event since there's no HR feed, "
    "extending the additive-only permission-override model without creating conflicts, safe role-inheritance cascading, "
    "scalability of the sweep/cache-invalidation path, a real rollback strategy for a failed mover/leaver step, "
    "HIPAA-defensible audit evidence, and a future AI enhancement -- anomaly detection on lifecycle patterns. Frame "
    "these as solvable, sequenced engineering risks, not blockers.")

# ---------------------------------------------------------------- Slide 4 --
# UC2 - Pain Points
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
add_header(s, "Use Case 2 · Pain Points", "Intelligent Ticket Lifecycle Automation", badge_text="RULES-BASED · NO AI COST")

left_x = CL
left_w = Inches(6.9)
add_text(s, left_x, Inches(1.42), left_w, Inches(0.3), "Current Manual Workflow", size=13.5, bold=True, color=NAVY)

pain_items2 = [
    ("MTC", "Manual Ticket Creation",
     "Every inbound email, call, or SMS is read and keyed into a ticket entirely by hand."),
    ("DUP", "Duplicate Work",
     "No detection of an existing thread — the same issue is often logged as two separate tickets."),
    ("ENT", "Manual Data Entry",
     "Client, claim number, category, and priority are typed in from memory or re-reading the message."),
    ("PRI", "Inconsistent Prioritization",
     "Priority depends on whoever triages it, not a consistent, repeatable signal."),
    ("FUP", "Follow-Up Burden",
     "Clients re-send or escalate simply because they received no visibility into ticket status."),
]
row_y = 1.84
for code, title, desc in pain_items2:
    pain_row(s, left_x, Inches(row_y), left_w, title, desc, code)
    row_y += 0.95

right_x = Inches(7.78)
right_w = Inches(5.0)
add_text(s, right_x, Inches(1.42), right_w, Inches(0.3), "Business Impact", size=13.5, bold=True, color=NAVY)
impacts2 = [
    ("Operational Delay", "Manual triage adds minutes-to-hours of latency before a ticket is even assigned.", RED, RED_BG),
    ("Customer Experience", "No Amazon-style tracking — clients are left guessing where their request stands.", AMBER, AMBER_BG),
    ("Data Quality", "Hand-typed intake data drives inconsistent reporting and misrouted tickets.", ACCENT, LIGHT_BLUE),
    ("Cost of Rework", "Duplicate handling and misassignment create avoidable rework for agents and supervisors.", NAVY, LIGHT_BLUE2),
]
cy = 1.84
for title, desc, bar, bg in impacts2:
    callout_card(s, right_x, Inches(cy), right_w, Inches(1.06), title, desc, bar, bg)
    cy += 1.19

add_footer(s, 5, TOTAL_SLIDES)
set_notes(s, 4,
    "Every inbound email/call/SMS is read and keyed in by hand today, with no duplicate detection -- the same issue "
    "is sometimes logged twice. Client, claim number, category, and priority are typed from memory. Priority is "
    "whoever-triages-it, not a consistent signal. Clients get no visibility into status, so they re-send or escalate "
    "simply from silence. Business impact: added triage latency, poor customer experience versus modern e-commerce-style "
    "tracking expectations, inconsistent reporting from hand-typed data, and avoidable rework from duplicates/misassignment.")

# ---------------------------------------------------------------- Slide 5 --
# UC2 - Solution Architecture
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
add_header(s, "Use Case 2 · Solution Architecture", "Intelligent Ticket Lifecycle Automation", title_size=23, badge_text="RULES-BASED · NO AI COST")

CX = Inches(3.15)
CWID = Inches(6.7)
RAIL_X = Inches(10.05)
RAIL_W = Inches(2.73)
LRAIL_X = Inches(0.55)
LRAIL_W = Inches(2.4)

# R1 - intake
process_box(s, CX, Inches(1.42), CWID, Inches(0.58), "Communication Intake & Gateway",
            "Email · SMS · Phone · Portal · Manual entry — normalized into one clean message object.",
            "CHN", icon_fill=NAVY, badge_size=Inches(0.32), title_size=10, desc_size=7.6)
add_connector(s, CX + CWID / 2, Inches(2.00), CX + CWID / 2, Inches(2.09), color=ACCENT)
add_text(s, CX + CWID / 2 + Inches(0.08), Inches(1.985), Inches(2.0), Inches(0.18), "Message Received", size=7, italic=True, color=TEXT_GRAY)

# R2 - rules-based extraction + deterministic duplicate matching
bw = Inches((6.7 - 0.2) / 2)
process_box(s, CX, Inches(2.10), bw, Inches(0.6), "Structured Parsing & Rules Engine",
            "Extracts client, claim #, patient, category & priority via regex, keyword rules and sender lookup — no model calls.", "RX", badge_size=Inches(0.3), title_size=9.3, desc_size=7.2)
process_box(s, CX + bw + Inches(0.2), Inches(2.10), bw, Inches(0.6), "Deterministic Matching Engine",
            "Matches on claim #, thread ID, and fuzzy text similarity — plain algorithms, no AI.", "DUP", badge_size=Inches(0.3), title_size=9.3, desc_size=7.2)
add_connector(s, CX + CWID / 2, Inches(2.70), CX + CWID / 2, Inches(2.79), color=ACCENT)
add_text(s, CX + CWID / 2 + Inches(0.08), Inches(2.685), Inches(1.6), Inches(0.18), "Match Signal", size=7, italic=True, color=TEXT_GRAY)

# R3 - decision diamond
dw = Inches(3.0)
dh = Inches(0.64)
dx = CX + (CWID - dw) / 2
decision_diamond(s, dx, Inches(2.80), dw, dh, "Attach or Create — Recommendation Only", size=8.5)
add_connector(s, dx + dw / 2, Inches(3.44), dx + dw / 2, Inches(3.53), color=ACCENT)
add_text(s, dx + dw / 2 + Inches(0.08), Inches(3.435), Inches(1.8), Inches(0.18), "Recommendation", size=7, italic=True, color=TEXT_GRAY)

# left dashed connector - manager override callout
add_connector(s, dx, Inches(2.80) + dh / 2, LRAIL_X + LRAIL_W, Inches(2.80) + dh / 2, color=ACCENT, dashed=True, width=1.1)
add_text(s, LRAIL_X + LRAIL_W - Inches(1.0), Inches(2.80) + dh / 2 - Inches(0.22), Inches(1.0), Inches(0.18),
         "Editable", size=7, italic=True, color=ACCENT, align=PP_ALIGN.RIGHT)

# R4 - Account Manager Review (highlighted)
rev = add_rect(s, CX, Inches(3.55), CWID, Inches(0.62), fill=ACCENT, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1, shadow=True)
add_icon_badge(s, CX + Inches(0.1), Inches(3.55) + Inches(0.09), Inches(0.34), "MGR", fill=WHITE, text_color=ACCENT, font_size=8)
add_text(s, CX + Inches(0.54), Inches(3.55) + Inches(0.08), Inches(4.6), Inches(0.26),
         "Account Manager Review — Human-in-the-Loop", size=10.5, bold=True, color=WHITE)
add_text(s, CX + Inches(0.54), Inches(3.55) + Inches(0.32), CWID - Inches(2.2), Inches(0.26),
         "Manager reviews the prefilled ticket, edits if needed, and makes the final Create / Attach decision.",
         size=7.6, color=WHITE)
badge = add_rect(s, CX + CWID - Inches(1.55), Inches(3.55) + Inches(0.1), Inches(1.4), Inches(0.24), fill=NAVY,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
fill_shape_text(badge, "FINAL AUTHORITY", size=6.8, bold=True, color=WHITE)
add_connector(s, CX + CWID / 2, Inches(4.17), CX + CWID / 2, Inches(4.25), color=ACCENT)
add_text(s, CX + CWID / 2 + Inches(0.08), Inches(4.155), Inches(2.0), Inches(0.18), "Manager Confirms", size=7, italic=True, color=TEXT_GRAY)

# R5 - ticket service
process_box(s, CX, Inches(4.26), CWID, Inches(0.5), "Ticket Service",
            "Creates or attaches the ticket, starts SLA, stores attachments and the audit entry.", "TKT", badge_size=Inches(0.3), title_size=9.5, desc_size=7.4)
add_connector(s, CX + CWID / 2, Inches(4.76), CX + CWID / 2, Inches(4.85), color=ACCENT)

# R6 - assignment / priority / sla
n = 3
gap = Inches(0.16)
cbw2 = Emu(int((CWID - gap * (n - 1)) / n))
cbx = CX
eng_items = [
    ("AS", "Assignment Engine", "best-fit agent"),
    ("PRI", "Priority Engine", "tier x keyword rules"),
    ("SLA", "SLA Engine", "starts the clock"),
]
for code, title, cap in eng_items:
    compact_box(s, cbx, Inches(4.86), cbw2, Inches(0.64), title, cap, code, title_size=8.3, caption_size=6.8)
    cbx = Emu(int(cbx + cbw2 + gap))
add_connector(s, CX + CWID / 2, Inches(5.50), CX + CWID / 2, Inches(5.59), color=ACCENT)

# R7 - status automation
process_box(s, CX, Inches(5.60), CWID, Inches(0.5), "Status Automation Engine",
            "Drives the Received → Closed tracking timeline and auto-closes after a configurable silence window.",
            "ST", badge_size=Inches(0.3), title_size=9.5, desc_size=7.3)
add_connector(s, CX + CWID / 2, Inches(6.10), CX + CWID / 2, Inches(6.19), color=ACCENT)

# R8 - outputs
bw = Inches((6.7 - 0.2) / 2)
process_box(s, CX, Inches(6.21), bw, Inches(0.5), "Customer Portal",
            "Amazon-style live status for the client.", "PT", icon_fill=NAVY, badge_size=Inches(0.28), title_size=9, desc_size=7)
process_box(s, CX + bw + Inches(0.2), Inches(6.21), bw, Inches(0.5), "Management Dashboard",
            "Load, SLA risk & rule-override rate for supervisors.", "MG", icon_fill=NAVY, badge_size=Inches(0.28), title_size=9, desc_size=7)

# Right rail
rail_items2 = [
    ("NT", "Notification Engine", "Fires a client / agent notification on every status transition."),
    ("AUD", "Audit Service", "Immutable record of every rule-based suggestion and every manager decision."),
    ("ANL", "Analytics", "Tracks rule-suggestion accept/override rate -- the tuning signal for the rules engine."),
]
rail_y2 = [2.10, 3.28, 4.46]
rail_h = Inches(0.98)
for (code, title, desc), ry in zip(rail_items2, rail_y2):
    process_box(s, RAIL_X, Inches(ry), RAIL_W, rail_h, title, desc, code, fill=LIGHT_BLUE2, border=ACCENT_SOFT,
                icon_fill=ACCENT, badge_size=Inches(0.3), title_size=9, desc_size=7.3)
for i in range(len(rail_y2) - 1):
    add_connector(s, RAIL_X + RAIL_W / 2, Inches(rail_y2[i]) + rail_h, RAIL_X + RAIL_W / 2, Inches(rail_y2[i + 1]), color=ACCENT_SOFT, width=1)
add_connector(s, CX + CWID, Inches(2.40), RAIL_X, Inches(rail_y2[0]) + rail_h / 2, color=ACCENT_SOFT, dashed=True, width=1, arrow=False)
add_connector(s, CX + CWID, Inches(3.86), RAIL_X, Inches(rail_y2[1]) + rail_h / 2, color=ACCENT_SOFT, dashed=True, width=1, arrow=False)
add_connector(s, CX + CWID, Inches(5.18), RAIL_X, Inches(rail_y2[2]) + rail_h / 2, color=ACCENT_SOFT, dashed=True, width=1, arrow=False)

# Left rail - manager override callout
callout_card(s, LRAIL_X, Inches(3.10), LRAIL_W, Inches(1.15),
             "Manager Override — Always Human-Controlled",
             "Every rule-based recommendation — duplicate match, category, priority, assignee — is editable before it's "
             "committed. Nothing is auto-created without a decision.",
             ACCENT, LIGHT_BLUE, title_size=9.5, desc_size=7.8)

add_footer(s, 6, TOTAL_SLIDES)
set_notes(s, 5,
    "Communication Channels feed a Gateway that normalizes intake, then a Structured Parsing & Rules Engine and a "
    "Deterministic Matching Engine run in parallel -- regex/keyword rules and sender lookup extract client/claim/patient/"
    "category/priority, and exact-ID plus fuzzy-text matching checks for a likely existing ticket. No AI/LLM call is made "
    "at all -- this keeps per-ticket cost near zero and every decision explainable by a plain rule, which matters for a "
    "HIPAA compliance review. This produces a recommendation only: Attach or Create. The critical design decision is the "
    "Account Manager Review stage -- highlighted in blue as Final Authority -- every rule-based suggestion is editable "
    "and nothing is auto-created without a human decision, shown again via the Manager Override callout on the left. Once "
    "confirmed, Ticket Service creates/attaches the ticket and starts SLA; Assignment, Priority, and SLA engines each "
    "make a rule-based recommendation the manager can override. Status Automation Engine drives the Amazon-style tracking "
    "timeline through to auto-close. Notification, Audit, and Analytics run as cross-cutting services on the right, and "
    "Analytics tracks how often each rule-based suggestion is accepted versus overridden -- the exact signal that would "
    "justify adding AI later for a specific underperforming field, instead of assuming it's needed upfront.")

# ---------------------------------------------------------------- Slide 6 --
# UC2 - Implementation Challenges
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
add_header(s, "Use Case 2 · Implementation Challenges", "Intelligent Ticket Lifecycle Automation", badge_text="RULES-BASED · NO AI COST")

challenges_uc2 = [
    ("RUL", "Rule Coverage & Maintenance", "A new client email format or claim-number pattern needs a new rule — a lightweight governance process keeps rules current, with no ML retraining cost."),
    ("DUP", "Duplicate Detection", "Overly aggressive fuzzy-matching thresholds risk false attach recommendations — must stay a suggestion, not an auto-merge."),
    ("OVR", "Manager Override", "The workflow must stay fast even though every rule-based suggestion requires a human confirmation step."),
    ("NTF", "Notification Reliability", "Every status transition must reliably notify the client — a missed notification defeats the tracking model."),
    ("INT", "Integrations", "Email / SMS / phone intake channels each need a resilient, idempotent ingestion connector."),
    ("SCL", "Scalability", "Rule evaluation and duplicate search must stay fast as volume grows — rules are far cheaper to scale than model inference."),
    ("CMP", "Compliance", "PHI-adjacent extracted fields must be handled under the same access controls as the ticket itself."),
    ("AUD", "Audit", "Every rule-based suggestion and every manager decision needs its own immutable audit entry."),
    ("SEC", "Security", "Rule evaluation must never bypass existing category / client visibility scoping."),
    ("DAT", "Data-Driven AI Trigger", "Track override rate per field; only add AI/LLM extraction for the specific fields where rules underperform — cost added only where the data justifies it."),
]
cols = 5
gap = Inches(0.16)
card_w = Emu(int((CW - gap * (cols - 1)) / cols))
card_h = Inches(2.35)
y0 = 1.55
row_gap = 0.35
for i, (code, title, desc) in enumerate(challenges_uc2):
    r, c = divmod(i, cols)
    x = Emu(int(CL + c * (card_w + gap)))
    y = Inches(y0 + r * (2.35 + row_gap))
    challenge_card(s, x, y, card_w, card_h, title, desc, code, title_size=9.7, desc_size=7.6)

add_footer(s, 7, TOTAL_SLIDES)
set_notes(s, 6,
    "Ten implementation risks to flag: rule coverage needs a lightweight governance process as new email formats and "
    "claim-number patterns appear -- but with no ML retraining cost; duplicate detection must stay a suggestion, never an "
    "auto-merge; the manager-override step must not slow the workflow down; notification reliability is core to the "
    "tracking experience; every intake channel needs a resilient, idempotent connector; rule evaluation must scale with "
    "volume, which is inherently cheaper than model inference; PHI-adjacent extracted fields inherit the ticket's own "
    "access controls; every rule-based suggestion and manager decision needs its own audit entry; rule evaluation must "
    "never bypass existing visibility scoping; and AI is a deliberate, data-driven future option -- add it only for the "
    "specific fields where the override-rate data shows rules underperforming, not as an upfront assumption.")

prs.save("C:/Unified-ticket-management-system/_probe.pptx")
print("all stages ok")
